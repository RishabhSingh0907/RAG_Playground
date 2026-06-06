"""
Multi-format document parsers for the ingestion pipeline.

Supports PDF, DOCX, TXT, CSV, PPT with metadata extraction and layout preservation.
Uses LlamaIndex and related libraries for robust parsing.
"""

import os
from pathlib import Path
from typing import Optional, List, Dict, Any
from datetime import datetime
import uuid

from src.ingestion.models import Document, DocumentMetadata
from src.utils.logger import get_logger


logger = get_logger(__name__)


class BaseParser:
    """Base class for document parsers."""
    
    supported_format: str
    
    def __init__(self):
        """Initialize parser."""
        pass
    
    def can_parse(self, file_path: str) -> bool:
        """Check if parser can handle this file format."""
        ext = Path(file_path).suffix.lower().lstrip(".")
        return ext == self.supported_format
    
    def parse(self, file_path: str) -> Optional[Document]:
        """
        Parse document and return Document object.
        
        Args:
            file_path: Path to file to parse
        
        Returns:
            Document object or None if parsing fails
        
        Raises:
            IOError: If file cannot be read
            ValueError: If format is not supported
        """
        raise NotImplementedError("Subclasses must implement parse()")
    
    def extract_metadata(self, file_path: str) -> DocumentMetadata:
        """Extract metadata from file."""
        file_obj = Path(file_path)
        
        return DocumentMetadata(
            file_path=str(file_path),
            file_name=file_obj.name,
            file_format=self.supported_format,
            file_size=file_obj.stat().st_size,
            created_at=datetime.fromtimestamp(file_obj.stat().st_ctime),
        )


class TextParser(BaseParser):
    """Parser for plain text files."""
    
    supported_format = "txt"
    
    def parse(self, file_path: str) -> Optional[Document]:
        """Parse text file."""
        try:
            logger.debug(f"Parsing text file: {file_path}")
            
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            
            if not content.strip():
                logger.warning(f"Empty text file: {file_path}")
                return None
            
            metadata = self.extract_metadata(file_path)
            document_id = str(uuid.uuid4())
            
            document = Document(
                document_id=document_id,
                content=content,
                metadata=metadata,
            )
            
            logger.info(f"Successfully parsed text file", file_path=file_path, doc_id=document_id)
            return document
            
        except IOError as e:
            logger.error(f"Failed to read text file: {file_path}", exc_info=True, error=str(e))
            return None
        except Exception as e:
            logger.error(f"Error parsing text file: {file_path}", exc_info=True, error=str(e))
            return None


class PDFParser(BaseParser):
    """Parser for PDF files using PyMuPDF."""
    
    supported_format = "pdf"
    
    def parse(self, file_path: str) -> Optional[Document]:
        """
        Parse PDF file using PyMuPDF (fitz).
        
        Requires: pip install pymupdf (or: pip install PyMuPDF)
        """
        try:
            logger.debug(f"Parsing PDF file: {file_path}")
            
            # Check file exists
            if not Path(file_path).exists():
                raise IOError(f"PDF file not found: {file_path}")
            
            # Parse with PyMuPDF
            content = self._parse_with_pymupdf(file_path)
            
            if content is None or not content.strip():
                logger.warning(f"No content extracted from PDF: {file_path}")
                return None
            
            metadata = self.extract_metadata(file_path)
            document_id = str(uuid.uuid4())
            
            document = Document(
                document_id=document_id,
                content=content,
                metadata=metadata,
            )
            
            logger.info(f"Successfully parsed PDF file", file_path=file_path, doc_id=document_id)
            return document
            
        except Exception as e:
            logger.error(f"Error parsing PDF file: {file_path}", exc_info=True, error=str(e))
            return None
    
    def _parse_with_pymupdf(self, file_path: str) -> Optional[str]:
        """Parse PDF using PyMuPDF library."""
        try:
            import fitz  # PyMuPDF
            
            logger.debug("Using PyMuPDF for PDF parsing")
            doc = fitz.open(file_path)
            content_parts = []
            
            for page_num, page in enumerate(doc):
                text = page.get_text()
                if text:
                    content_parts.append(f"--- Page {page_num + 1} ---\n{text}")
            
            content = "\n\n".join(content_parts)
            return content if content.strip() else None
            
        except ImportError:
            logger.error("PyMuPDF not installed. Install with: pip install pymupdf")
            return None
        except Exception as e:
            logger.error(f"PyMuPDF parsing failed: {str(e)}", exc_info=True, error=str(e))
            return None


class DOCXParser(BaseParser):
    """Parser for DOCX (Word) files using DoclingLoader."""
    
    supported_format = "docx"
    
    def parse(self, file_path: str) -> Optional[Document]:
        """
        Parse DOCX file.
        
        Primary: pip install langchain-docling (for langchain_docling.loader.DoclingLoader)
        Fallback: pip install python-docx
        """
        try:
            logger.debug(f"Parsing DOCX file: {file_path}")
            
            if not Path(file_path).exists():
                raise IOError(f"DOCX file not found: {file_path}")
            
            # Try DoclingLoader first (primary method)
            content = self._parse_with_docling(file_path)
            
            # Fallback to python-docx if DoclingLoader fails
            if content is None:
                logger.debug("DoclingLoader parsing failed, attempting python-docx...")
                content = self._parse_with_python_docx(file_path)
            
            if not content or not content.strip():
                logger.warning(f"Empty content from DOCX: {file_path}")
                return None
            
            metadata = self.extract_metadata(file_path)
            document_id = str(uuid.uuid4())
            
            document = Document(
                document_id=document_id,
                content=content,
                metadata=metadata,
            )
            
            logger.info(f"Successfully parsed DOCX file", file_path=file_path, doc_id=document_id)
            return document
            
        except Exception as e:
            logger.error(f"Error parsing DOCX file: {file_path}", exc_info=True, error=str(e))
            return None
    
    def _parse_with_docling(self, file_path: str) -> Optional[str]:
        """Parse DOCX using DoclingLoader (primary method)."""
        try:
            from langchain_docling.loader import DoclingLoader
            
            logger.debug("Using DoclingLoader for DOCX parsing")
            loader = DoclingLoader(file_path=file_path)
            documents = loader.load()
            
            if not documents:
                return None
            
            # Extract content from loaded documents
            content_parts = []
            for doc in documents:
                if hasattr(doc, 'page_content'):
                    content_parts.append(doc.page_content)
                elif hasattr(doc, 'content'):
                    content_parts.append(doc.content)
                else:
                    content_parts.append(str(doc))
            
            content = "\n".join(content_parts)
            return content if content.strip() else None
            
        except ImportError:
            logger.debug("DoclingLoader not installed. Install with: pip install langchain-docling")
            return None
        except Exception as e:
            logger.debug(f"DoclingLoader parsing failed: {str(e)}")
            return None
    
    def _parse_with_python_docx(self, file_path: str) -> Optional[str]:
        """Parse DOCX using python-docx library (fallback method)."""
        try:
            from docx import Document as DocxDocument
            
            logger.debug("Using python-docx for DOCX parsing")
            doc = DocxDocument(file_path)
            
            # Extract text from paragraphs and tables
            content_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    content_parts.append(para.text)
            
            # Extract table content
            for table in doc.tables:
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    content_parts.append(" | ".join(row_data))
            
            content = "\n".join(content_parts)
            return content if content.strip() else None
            
        except ImportError:
            logger.debug("python-docx not installed. Install with: pip install python-docx")
            return None
        except Exception as e:
            logger.debug(f"python-docx parsing failed: {str(e)}")
            return None


class CSVParser(BaseParser):
    """Parser for CSV files."""
    
    supported_format = "csv"
    
    def parse(self, file_path: str) -> Optional[Document]:
        """
        Parse CSV file.
        
        Requires: pip install pandas
        """
        try:
            logger.debug(f"Parsing CSV file: {file_path}")
            
            if not Path(file_path).exists():
                raise IOError(f"CSV file not found: {file_path}")
            
            try:
                import pandas as pd
            except ImportError:
                logger.error("pandas not installed. Install with: pip install pandas")
                return None
            
            # Parse CSV
            df = pd.read_csv(file_path)
            
            if df.empty:
                logger.warning(f"Empty CSV file: {file_path}")
                return None
            
            # Convert to text representation
            content = df.to_string()
            
            metadata = self.extract_metadata(file_path)
            document_id = str(uuid.uuid4())
            
            document = Document(
                document_id=document_id,
                content=content,
                metadata=metadata,
            )
            
            logger.info(f"Successfully parsed CSV file", file_path=file_path, doc_id=document_id, rows=len(df))
            return document
            
        except ImportError as e:
            logger.error(f"Missing dependency for CSV parsing: {str(e)}", exc_info=True)
            return None
        except Exception as e:
            logger.error(f"Error parsing CSV file: {file_path}", exc_info=True, error=str(e))
            return None


class PPTParser(BaseParser):
    """Parser for PowerPoint files."""
    
    supported_format = "ppt"
    
    def parse(self, file_path: str) -> Optional[Document]:
        """
        Parse PowerPoint file (.ppt or .pptx).
        
        Requires: pip install python-pptx
        """
        try:
            logger.debug(f"Parsing PowerPoint file: {file_path}")
            
            if not Path(file_path).exists():
                raise IOError(f"PowerPoint file not found: {file_path}")
            
            # Handle both .ppt and .pptx (we treat both as "ppt" format)
            file_ext = Path(file_path).suffix.lower()
            if file_ext == ".pptx":
                try:
                    from pptx import Presentation
                except ImportError:
                    logger.error("python-pptx not installed. Install with: pip install python-pptx")
                    return None
                
                # Parse PPTX
                prs = Presentation(file_path)
                content_parts = []
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    content_parts.append(f"--- Slide {slide_num} ---")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            content_parts.append(shape.text)
                
                content = "\n".join(content_parts)
            else:
                logger.warning(f"Older .ppt format not fully supported: {file_path}. Please convert to .pptx")
                return None
            
            if not content.strip():
                logger.warning(f"Empty content from PowerPoint: {file_path}")
                return None
            
            metadata = self.extract_metadata(file_path)
            document_id = str(uuid.uuid4())
            
            document = Document(
                document_id=document_id,
                content=content,
                metadata=metadata,
            )
            
            logger.info(f"Successfully parsed PowerPoint file", file_path=file_path, doc_id=document_id)
            return document
            
        except ImportError as e:
            logger.error(f"Missing dependency for PowerPoint parsing: {str(e)}", exc_info=True)
            return None
        except Exception as e:
            logger.error(f"Error parsing PowerPoint file: {file_path}", exc_info=True, error=str(e))
            return None


class ParserRegistry:
    """Registry of available parsers for multimodal document processing."""
    
    def __init__(self):
        """Initialize parser registry with all supported parsers."""
        self.parsers: Dict[str, BaseParser] = {
            "txt": TextParser(),
            "pdf": PDFParser(),
            "docx": DOCXParser(),
            "csv": CSVParser(),
            "ppt": PPTParser(),  # Includes .pptx
            "pptx": PPTParser(),
        }
        logger.info(f"ParserRegistry initialized with {len(self.parsers)} parsers")
    
    def get_parser(self, file_path: str) -> Optional[BaseParser]:
        """Get appropriate parser for file."""
        ext = Path(file_path).suffix.lower().lstrip(".")
        
        if ext not in self.parsers:
            logger.warning(f"No parser available for format: {ext}")
            return None
        
        return self.parsers[ext]
    
    def parse_file(self, file_path: str) -> Optional[Document]:
        """
        Parse file using appropriate parser.
        
        Args:
            file_path: Path to file to parse
        
        Returns:
            Document object or None if parsing fails
        """
        parser = self.get_parser(file_path)
        
        if not parser:
            logger.error(f"No parser available for file: {file_path}")
            return None
        
        return parser.parse(file_path)
    
    def register_parser(self, format_name: str, parser: BaseParser) -> None:
        """
        Register a custom parser.
        
        Args:
            format_name: File format name (e.g., "pdf", "docx")
            parser: Parser instance
        """
        self.parsers[format_name.lower()] = parser
        logger.info(f"Registered custom parser for format: {format_name}")
